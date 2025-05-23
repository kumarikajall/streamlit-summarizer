import os
import PyPDF2
import docx
from pptx import Presentation
from transformers import PegasusTokenizer, PegasusForConditionalGeneration
import torch

class PegasusFileSummarizer:
    def __init__(self, model_name="google/pegasus-xsum"):
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        self.tokenizer = PegasusTokenizer.from_pretrained(model_name)
        self.model = PegasusForConditionalGeneration.from_pretrained(model_name).to(self.device)

    def extract_text(self, file_path):
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        if ext == '.pdf':
            return self._extract_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            return self._extract_docx(file_path)
        elif ext in ['.pptx', '.ppt']:
            return self._extract_pptx(file_path)
        elif ext == '.txt':
            return self._extract_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    def _extract_pdf(self, file_path):
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            return '\n'.join([page.extract_text() or '' for page in reader.pages])

    def _extract_docx(self, file_path):
        doc = docx.Document(file_path)
        return '\n'.join([para.text for para in doc.paragraphs])

    def _extract_pptx(self, file_path):
        prs = Presentation(file_path)
        return '\n'.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')])

    def _extract_txt(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()

    def summarize(self, text, max_length=150, min_length=40):
        inputs = self.tokenizer(text, return_tensors="pt", truncation=True, max_length=1024).to(self.device)
        summary_ids = self.model.generate(
            inputs["input_ids"],
            max_length=max_length,
            min_length=min_length,
            length_penalty=2.0,
            num_beams=4,
            early_stopping=True
        )
        return self.tokenizer.decode(summary_ids[0], skip_special_tokens=True)

    def summarize_file(self, file_path, max_length=150, min_length=40):
        text = self.extract_text(file_path)
        return self.summarize(text, max_length, min_length)
