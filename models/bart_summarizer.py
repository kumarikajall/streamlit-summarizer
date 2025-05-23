import os
import docx
from pptx import Presentation
import PyPDF2
from transformers import BartTokenizer, BartForConditionalGeneration
import torch

class BARTFileSummarizer:
    def __init__(self):
        self.tokenizer = BartTokenizer.from_pretrained('facebook/bart-large-cnn')
        self.model = BartForConditionalGeneration.from_pretrained('facebook/bart-large-cnn')
        self.device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
        self.model = self.model.to(self.device)
        self.max_chunk_len = 1024

    def extract_text(self, file_path):
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        if ext == '.pdf':
            return self._extract_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            return self._extract_docx(file_path)
        elif ext in ['.pptx', '.ppt']:
            return self._extract_pptx(file_path)
        elif ext == '.txt':
            return self._extract_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    def _extract_pdf(self, file_path):
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            return '\n'.join([page.extract_text() or '' for page in reader.pages])

    def _extract_docx(self, file_path):
        doc = docx.Document(file_path)
        return '\n'.join([para.text for para in doc.paragraphs])

    def _extract_pptx(self, file_path):
        prs = Presentation(file_path)
        runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    runs.append(shape.text)
        return '\n'.join(runs)

    def _extract_txt(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()

    def summarize(self, text, max_length=150, min_length=40):
        chunks = [text[i:i+1000] for i in range(0, len(text), 1000)]
        inputs = self.tokenizer.batch_encode_plus(
            chunks, max_length=self.max_chunk_len, truncation=True,
            padding='longest', return_tensors='pt'
        )
        inputs = {k: v.to(self.device) for k, v in inputs.items()}
        outputs = self.model.generate(
            input_ids=inputs['input_ids'],
            attention_mask=inputs['attention_mask'],
            max_length=max_length,
            min_length=min_length,
            length_penalty=2.0,
            num_beams=4,
            early_stopping=True
        )
        return ' '.join(self.tokenizer.decode(o, skip_special_tokens=True) for o in outputs)

    def summarize_file(self, file_path, max_length=150, min_length=40):
        text = self.extract_text(file_path)
        return self.summarize(text, max_length, min_length)
